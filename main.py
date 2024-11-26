from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches
from PIL import Image

# Open file PowerPoint
prs = Presentation('path/to/your/presentation.pptx')

# Choose slide number
slide = prs.slides[0]

# Save slides
for index, shape in enumerate(slide.shapes):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        image = shape.image
        image_bytes = image.blob
        with open(f'slide_image_{index}.png', 'wb') as image_file:
            image_file.write(image_bytes)